from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def add_icon_shape(slide, icon_char, left, top, width=Inches(0.6), height=Inches(0.6), font_size=Pt(28)):
    """添加图标文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon_char
    p.font.size = font_size
    p.alignment = PP_ALIGN.CENTER
    return txBox


def add_title_slide(prs, title, subtitle, speaker="", date=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(25, 55, 95)
    shape.line.fill.background()

    # Decorative circles
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(0.5), Inches(1), Inches(1))
    c1.fill.solid()
    c1.fill.fore_color.rgb = RGBColor(45, 85, 130)
    c1.line.fill.background()

    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.3), Inches(5.8), Inches(0.8), Inches(0.8))
    c2.fill.solid()
    c2.fill.fore_color.rgb = RGBColor(45, 85, 130)
    c2.line.fill.background()

    # Icons row
    icons = ["⚙️", "🔧", "📐", "💪"]
    for i, icon in enumerate(icons):
        add_icon_shape(slide, icon, Inches(2 + i * 1.5), Inches(1.8), font_size=Pt(32))

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(200, 220, 255)
    p2.alignment = PP_ALIGN.CENTER

    # Speaker & Date
    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = f"{speaker}  |  {date}" if speaker and date else ""
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(180, 200, 230)
    p3.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets, section="", icon=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Light background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
    shape.line.fill.background()

    # Section header bar
    if section:
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(25, 55, 95)
        header.line.fill.background()

        sec_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
        sec_tf = sec_box.text_frame
        sec_p = sec_tf.paragraphs[0]
        sec_p.text = f"{icon} {section}" if icon else section
        sec_p.font.size = Pt(14)
        sec_p.font.color.rgb = RGBColor(180, 200, 230)

    # Title with icon
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{icon} {title}" if icon else title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 55, 95)

    # Bullets
    bullet_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(8.6), Inches(4.5))
    tf2 = bullet_box.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"▶  {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(14)

    return slide


def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets, section="", icon=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Light background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
    shape.line.fill.background()

    # Section header bar
    if section:
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(25, 55, 95)
        header.line.fill.background()

        sec_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
        sec_tf = sec_box.text_frame
        sec_p = sec_tf.paragraphs[0]
        sec_p.text = f"{icon} {section}" if icon else section
        sec_p.font.size = Pt(14)
        sec_p.font.color.rgb = RGBColor(180, 200, 230)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{icon} {title}" if icon else title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 55, 95)

    # Left column title with icon
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(4.3), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"📂 {left_title}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 55, 95)

    # Left column bullets
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(4.3), Inches(3.5))
    tf2 = left_box.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(10)

    # Right column title with icon
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.1), Inches(4.3), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"📂 {right_title}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 55, 95)

    # Right column bullets
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.7), Inches(4.3), Inches(3.5))
    tf3 = right_box.text_frame
    tf3.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(10)

    return slide


def add_closing_slide(prs, quote, thank_you="谢谢"):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Dark background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(25, 55, 95)
    shape.line.fill.background()

    # Decorative elements
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(5.5), Inches(1.5), Inches(1.5))
    c1.fill.solid()
    c1.fill.fore_color.rgb = RGBColor(45, 85, 130)
    c1.line.fill.background()

    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(0.5), Inches(1.2), Inches(1.2))
    c2.fill.solid()
    c2.fill.fore_color.rgb = RGBColor(45, 85, 130)
    c2.line.fill.background()

    # Icons
    add_icon_shape(slide, "🎓", Inches(4.3), Inches(1.2), font_size=Pt(40))

    # Quote
    quote_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f'"{quote}"'
    p.font.size = Pt(28)
    p.font.italic = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Thank you
    ty_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    tf2 = ty_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"🙏 {thank_you}"
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(200, 220, 255)
    p2.alignment = PP_ALIGN.CENTER

    # Q&A
    qa_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.5))
    tf3 = qa_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "❓ Q & A"
    p3.font.size = Pt(20)
    p3.font.color.rgb = RGBColor(180, 200, 230)
    p3.alignment = PP_ALIGN.CENTER

    return slide


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Cover
    add_title_slide(
        prs,
        "力学之道，成长之路",
        "大学生力学学习核心建议 + 人生经验之谈 🏆",
        "力学研究者", "2026年4月"
    )

    # Slide 2: 力学学习的底层逻辑
    add_content_slide(
        prs,
        "力学不是\"做题\"，而是构建\"思维模型\" 🧠",
        [
            "力学 = 数学工具 ⚙️ + 物理直觉 🔬 + 工程抽象 🏗️",
            "核心能力：从实际问题中提炼力学模型（这才是大佬的标配！）",
            "思维闭环：实际问题 → 简化假设 → 力学模型 → 数学方程 → 结果验证 → 工程判断",
            "敲黑板：理解\"为什么\"比死记\"怎么做\"重要100倍！"
        ],
        section="第一部分：如何高效学习力学专业",
        icon="📚"
    )

    # Slide 3: 基础不牢，地动山摇
    add_content_slide(
        prs,
        "基础不牢？地动山摇！🚧",
        [
            "理论力学（30%）：矢量分析 + 运动描述 = 你的肌肉记忆 💪",
            "材料力学（30%）：内力图 + 应力应变 = 直觉养成 📊",
            "数学物理方法（25%）：偏微分方程 + 矩阵论 = 躲不掉的空气 🌬️",
            "编程计算（15%）：MATLAB/Python = 打怪升级的必备武器 🎮"
        ],
        section="第一部分：如何高效学习力学专业",
        icon="🏗️"
    )

    # Slide 4: 建立力学可视化直觉
    add_two_column_slide(
        prs,
        "把公式\"画\"在脑子里 🎨",
        "工具与技能",
        [
            "MATLAB/Python 做模拟（代码比笔头快！）",
            "看懂应力云图、模态振型图、流线图",
            "用代码复现教材经典案例（卷死同学😎）"
        ],
        "实践建议",
        [
            "悬臂梁弯曲应力分布（红蓝渐变美如画）",
            "振动力学模态分析（看着振型睡不着的那种）",
            "CFD流线可视化（流体也能这么优雅）"
        ],
        section="第一部分：如何高效学习力学专业",
        icon="📈"
    )

    # Slide 5: 力学人的三把刷子 + 学习英语
    add_two_column_slide(
        prs,
        "\"三把刷子\" + 英语是特权 🌐",
        "三门硬核技能",
        [
            "有限元思维：理解单元、节点、边界条件的物理意义（不是点软件那么简单！）",
            "实验报告撰写：数据处理 + 误差分析 = 工程师的信誉背书",
            "文献检索：追踪顶刊 JMPS、IJSS（发论文的必经之路）"
        ],
        "英语让你开挂",
        [
            "力学前沿论文 99% 为英文（中文文献？大海捞针！）",
            "精读专业词汇，搞定英文推导（和老外谈笑风生）",
            "尝试写英文学术摘要（从此走上人生巅峰🚀）"
        ],
        section="第一部分 & 第二部分",
        icon="⚡"
    )

    # Slide 6: AI与复利思维
    add_two_column_slide(
        prs,
        "AI是你的免费助教 + 定投是时间的朋友 💰",
        "AI助力力学学习",
        [
            "解释复杂的本构方程推导逻辑（24小时在线的私教）",
            "快速生成 MATLAB 后处理绘图代码（摸鱼神器！）",
            "润色英文学术邮件与摘要（告别Chinglish）",
            "效率提升可达 3 倍以上（卷王养成的秘密）"
        ],
        "定投：对抗焦虑的良药",
        [
            "定投不仅是理财，更是心态（知识、技能双管齐下）",
            "每月存10%生活费买宽基指数（博士毕业时你会感谢我）",
            "复利拐点在5-7年后（和时间做朋友👊）"
        ],
        section="第二部分：我的人生经验之谈",
        icon="🤖"
    )

    # Slide 7: 坚持健身
    add_content_slide(
        prs,
        "力学研究拼到最后，拼的是体力 🏋️",
        [
            "力学是长周期脑力劳动，好的心肺功能是高强度计算的后盾",
            "建议：每周 3 次力量训练 + 2 次有氧（科研也要体能管理！）",
            "运动与专注力正相关：运动越多，连续专注时长越长（血赚不亏）",
            "身体是革命的本钱——也是你能卷过同龄人的终极武器 💪🔥"
        ],
        section="第二部分：我的人生经验之谈",
        icon="🏃"
    )

    # Slide 8: 结语
    add_closing_slide(
        prs,
        "用理性的力学思维解构世界，用自律的长期主义重构人生"
    )

    prs.save("力学之道_成长之路.pptx")
    print("PPT已生成：力学之道_成长之路.pptx")


if __name__ == "__main__":
    create_presentation()