import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors
BLUE = RGBColor(0, 82, 139)
ORANGE = RGBColor(255, 140, 0)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(100, 100, 100)
RED = RGBColor(220, 53, 69)
GREEN = RGBColor(40, 167, 69)
PURPLE = RGBColor(111, 66, 193)
CYAN = RGBColor(23, 162, 184)
YELLOW = RGBColor(255, 193, 7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_header(slide, title):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
    p = txBox.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

def add_icon_box(slide, emoji, text, x, y, size=0.6, bg_color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(x, y + Inches(0.1), Inches(size), Inches(size - 0.1))
    p = txBox.text_frame.paragraphs[0]
    p.text = emoji
    p.font.size = Pt(int(size * 28))
    p.alignment = PP_ALIGN.CENTER

def add_bullet_icon(slide, emoji, text, y, size=20, color=GRAY):
    txBox = slide.shapes.add_textbox(Inches(0.5), y, Inches(0.5), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = emoji
    p.font.size = Pt(size)
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(Inches(1), y, Inches(11.5), Inches(0.5))
    p2 = txBox2.text_frame.paragraphs[0]
    p2.text = text
    p2.font.size = Pt(size)
    p2.font.color.rgb = color

def add_bullet(slide, text, y, size=22, bold=False, color=GRAY):
    txBox = slide.shapes.add_textbox(Inches(0.6), y, Inches(12), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color

def add_section(slide, num, title):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"0{num}"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

# === SLIDE 1: Cover ===
slide = add_slide()
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
shape.fill.solid()
shape.fill.fore_color.rgb = BLUE
shape.line.fill.background()

# Decorative circles
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-1), Inches(4), Inches(4))
shape.fill.solid()
shape.fill.fore_color.rgb = ORANGE
shape.line.fill.background()
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(5), Inches(3), Inches(3))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0, 60, 100)
shape.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
p = txBox.text_frame.paragraphs[0]
p.text = "力学专业大一生存与发展指南"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(0.8))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "从基础夯实到多元发展的四年路线图"
p2.font.size = Pt(24)
p2.font.color.rgb = ORANGE
p2.alignment = PP_ALIGN.CENTER

# === SLIDE 2: 力学认知 + 三维出口 ===
slide = add_slide()
add_header(slide, "力学专业认知与就业方向")

# 左侧带图标
add_icon_box(slide, "🎯", "学科定位", Inches(0.5), Inches(1.3), 0.5, BLUE)
add_bullet(slide, '力学是"工程科学的基石"', Inches(1.3), 20)
add_bullet(slide, "上游：数学、物理", Inches(1.8), 17, color=GRAY)
add_bullet(slide, "下游：航空航天、机械、土木、金融", Inches(2.2), 17, color=GRAY)

# 三维出口标题
add_icon_box(slide, "🧭", "三维出口", Inches(0.5), Inches(3), 0.5, ORANGE)

# 右侧表格
headers = ["维度", "方向", "代表领域"]
rows = [
    ["📚 学术深造型", "固体力学/流体力学", "高校、国家重点实验室"],
    ["⚙️ 工程应用型", "CAE仿真/结构设计", "商飞、航天科技、汽车主机厂"],
    ["🌉 跨界融合型", "量化金融/数据科学", "投行、互联网、科技公司"]
]
table = slide.shapes.add_table(4, 3, Inches(5.5), Inches(1.3), Inches(7.3), Inches(2)).table
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE
    cell.text_frame.paragraphs[0].font.size = Pt(14)
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.color.rgb = WHITE
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = table.cell(r + 1, c)
        cell.text = val
        cell.text_frame.paragraphs[0].font.size = Pt(13)

# 底部破除误区
add_icon_box(slide, "💡", "", Inches(0.5), Inches(4.2), 0.4, YELLOW)
add_bullet(slide, '力学+编程=仿真工程师 | 力学+金融=量化分析师 | 力学+生物=生物力学研究员', Inches(4.3), 15, color=ORANGE)

# === SLIDE 3: 大一核心任务 ===
slide = add_slide()
add_header(slide, "大一核心任务：基础课突围")

# 2x2优先级矩阵 with emojis
# Box 1: 紧急且重要
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(5.8), Inches(1.6))
shape.fill.solid()
shape.fill.fore_color.rgb = RED
shape.line.fill.background()
add_icon_box(slide, "🔥", "", Inches(0.7), Inches(1.5), 0.4, WHITE)
txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(4.8), Inches(0.5))
p = txBox.text_frame.paragraphs[0]
p.text = "紧急且重要"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = WHITE
txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(2), Inches(4.8), Inches(0.8))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "高数、线代、大学物理（力学）"
p2.font.size = Pt(15)
p2.font.color.rgb = WHITE

# Box 2: 重要不紧急
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4), Inches(5.8), Inches(1.6))
shape.fill.solid()
shape.fill.fore_color.rgb = YELLOW
shape.line.fill.background()
add_icon_box(slide, "📝", "", Inches(7), Inches(1.5), 0.4, WHITE)
txBox = slide.shapes.add_textbox(Inches(7.5), Inches(1.5), Inches(4.8), Inches(0.5))
p = txBox.text_frame.paragraphs[0]
p.text = "重要不紧急"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = WHITE
txBox2 = slide.shapes.add_textbox(Inches(7.5), Inches(2), Inches(4.8), Inches(0.8))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "Python/MATLAB、工程制图"
p2.font.size = Pt(15)
p2.font.color.rgb = WHITE

# Box 3: 紧急不重要
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.3), Inches(5.8), Inches(1.6))
shape.fill.solid()
shape.fill.fore_color.rgb = CYAN
shape.line.fill.background()
add_icon_box(slide, "📋", "", Inches(0.7), Inches(3.4), 0.4, WHITE)
txBox = slide.shapes.add_textbox(Inches(1.2), Inches(3.4), Inches(4.8), Inches(0.5))
p = txBox.text_frame.paragraphs[0]
p.text = "紧急不重要"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = WHITE
txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(3.9), Inches(4.8), Inches(0.8))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "英语四六级（尽早过）"
p2.font.size = Pt(15)
p2.font.color.rgb = WHITE

# Box 4: 不紧急不重要
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(3.3), Inches(5.8), Inches(1.6))
shape.fill.solid()
shape.fill.fore_color.rgb = GRAY
shape.line.fill.background()
add_icon_box(slide, "⏸️", "", Inches(7), Inches(3.4), 0.4, WHITE)
txBox = slide.shapes.add_textbox(Inches(7.5), Inches(3.4), Inches(4.8), Inches(0.5))
p = txBox.text_frame.paragraphs[0]
p.text = "不紧急不重要"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = WHITE
txBox2 = slide.shapes.add_textbox(Inches(7.5), Inches(3.9), Inches(4.8), Inches(0.8))
p2 = txBox2.text_frame.paragraphs[0]
p2.text = "别浪费时间"
p2.font.size = Pt(15)
p2.font.color.rgb = WHITE

# 理论力学 section
add_icon_box(slide, "⚡", "理论力学", Inches(0.5), Inches(5.3), 0.5, GREEN)
add_bullet(slide, "静力学 → 运动学 → 动力学", Inches(5.9), 18, bold=True, color=BLUE)
add_bullet(slide, "核心：建立坐标系、受力分析、列方程的规范流程", Inches(6.4), 16, color=GRAY)
add_icon_box(slide, "🤔", "", Inches(0.5), Inches(6.7), 0.4, PURPLE)
add_bullet(slide, "关键：问自己——约束条件？自由度？能用能量法？", Inches(5.9), 16, color=ORANGE)

# === SLIDE 4: 工具与技能 ===
slide = add_slide()
add_header(slide, "工具与技能储备")

# 编程路径 with emoji timeline
add_icon_box(slide, "💻", "编程学习路径", Inches(0.5), Inches(1.3), 0.5, BLUE)

timeline_y = Inches(1.9)
labels = ["第1学期", "第2学期", "大二", "大三+"]
x_starts = [Inches(1.5), Inches(4.2), Inches(6.9), Inches(9.6)]
widths = [Inches(2.2), Inches(2.2), Inches(2.2), Inches(2.5)]
colors = [RGBColor(0, 123, 255), RGBColor(0, 123, 255), GREEN, PURPLE]
emojis = ["📖", "🔧", "🚀", "🎯"]

for i, (label, x, w, c, em) in enumerate(zip(labels, x_starts, widths, colors, emojis)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, timeline_y, w, Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = c
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(x, timeline_y + Inches(0.15), w, Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = f"{em} {label}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

contents = ["基础语法+NumPy", "SymPy/简单仿真", "有限元/FEM", "专业深化"]
for i, (x, content) in enumerate(zip(x_starts, contents)):
    txBox = slide.shapes.add_textbox(x, timeline_y + Inches(0.75), widths[i], Inches(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = content
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

# 工程软件表格
add_icon_box(slide, "🛠️", "工程软件", Inches(0.5), Inches(3.5), 0.5, ORANGE)
headers = ["软件类型", "代表软件", "用途", "时间"]
rows = [
    ["📐 三维建模", "SolidWorks/CATIA", "几何建模", "大二上"],
    ["🔬 有限元分析", "ANSYS/Abaqus", "应力/变形仿真", "大二下"],
    ["⚙️ 多体动力学", "Adams/Simpack", "机构运动分析", "大三"],
    ["🌊 计算流体力学", "Fluent/OpenFOAM", "流体仿真", "大三"]
]
table = slide.shapes.add_table(5, 4, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.2)).table
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE
    cell.text_frame.paragraphs[0].font.size = Pt(14)
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.color.rgb = WHITE
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = table.cell(r + 1, c)
        cell.text = val
        cell.text_frame.paragraphs[0].font.size = Pt(13)

# === SLIDE 5: 科研与竞赛 ===
slide = add_slide()
add_header(slide, "科研与竞赛启蒙")

# 科研策略
add_icon_box(slide, "🔬", '科研"轻接触"', Inches(0.5), Inches(1.3), 0.5, PURPLE)
add_bullet(slide, "了解科研是什么，不急于出成果", Inches(1.5), 17, color=GRAY)

actions = [
    ("👥", "旁听组会：联系老师，申请旁听研究生组会"),
    ("📖", "文献泛读：每周1篇中文综述，了解领域概况"),
    ("🎯", "方向选择：推导→固体力学 | 仿真→CFD | 实验→生物力学")
]
for i, (em, action) in enumerate(actions):
    add_icon_box(slide, em, "", Inches(0.5), Inches(2 + i * 0.65), 0.35, BLUE)
    add_bullet(slide, action, Inches(2.05 + i * 0.65), 16)

# 竞赛时间线
add_icon_box(slide, "🏆", "竞赛时间线", Inches(0.5), Inches(4.2), 0.5, YELLOW)
headers = ["时间", "竞赛", "准备要点"]
rows = [
    ["📅 大一寒假", "MCM/ICM数学建模", "组队+建模基础"],
    ["📅 大一下", "周培源力学竞赛", "理论+材料力学"],
    ["📅 大二", "结构设计竞赛", "团队+动手"],
    ["📅 持续", '"挑战杯"/"互联网+"', "科研转化"]
]
table = slide.shapes.add_table(5, 3, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2)).table
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = ORANGE
    cell.text_frame.paragraphs[0].font.size = Pt(14)
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.color.rgb = WHITE
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = table.cell(r + 1, c)
        cell.text = val
        cell.text_frame.paragraphs[0].font.size = Pt(13)

# === SLIDE 6: 职业方向 + 四年配方 ===
slide = add_slide()
add_header(slide, "职业方向与四年规划")

# 行业方向
add_icon_box(slide, "🗺️", "行业地图", Inches(0.5), Inches(1.3), 0.5, BLUE)
add_bullet(slide, "传统工程：航空航天（商飞/航天科工）| 汽车（比亚迪/蔚来）| 能源（风电/核电）", Inches(1.5), 15, color=GRAY)
add_bullet(slide, "新兴交叉：CAE仿真 | 机器人 | 生物力学 | 量化金融", Inches(2), 15, color=ORANGE)

# 四年配方表格
add_icon_box(slide, "📅", "四年配方", Inches(0.5), Inches(2.8), 0.5, GREEN)
headers = ["目标", "大一", "大二", "大三", "大四"]
rows = [
    ["🎓 保研", "高数线代满分", "全优+助研", "夏令营+推免", "毕设"],
    ["🌍 出国", "英语+GPA", "科研+GRE", "论文+套磁", "申请"],
    ["💼 就业", "编程基础", "实习1段", "实习2段", "入职"],
    ["🔄 跨专业", "辅修准备", "转专业", "转型项目", "过渡"]
]
table = slide.shapes.add_table(5, 5, Inches(0.3), Inches(3.4), Inches(12.7), Inches(2)).table
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE
    cell.text_frame.paragraphs[0].font.size = Pt(14)
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.color.rgb = WHITE
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = table.cell(r + 1, c)
        cell.text = val
        cell.text_frame.paragraphs[0].font.size = Pt(13)

# 大一Checklist
add_icon_box(slide, "✅", "大一Checklist", Inches(0.5), Inches(5.7), 0.5, ORANGE)
checklist = "高数/线代90+ | Python基础 | 旁听组会 | 过四级 | Office Hour认识老师 | 读《结构是什么》"
add_bullet(slide, checklist, Inches(6.2), 14, color=GRAY)

# === SLIDE 7: 四年时间轴 ===
slide = add_slide()
add_header(slide, "大学四年时间轴")

# 时间轴header
years = ["大一上", "大一下", "大二上", "大二下", "大三上", "大三下", "大四"]
xs = [Inches(1.8), Inches(3.4), Inches(5), Inches(6.6), Inches(8.2), Inches(9.8), Inches(11.4)]
for y, x in zip(years, xs):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), Inches(1.2), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(x, Inches(1.35), Inches(1.2), Inches(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = y
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# 时间轴线条
for x in xs:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.5), Inches(1.75), Inches(0.05), Inches(4.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(220, 220, 220)
    shape.line.fill.background()

# 横向条目 with emoji
entries = [
    ("📚 基础课", [(0, 2, BLUE)], "高数/线代/大物"),
    ("💻 编程", [(0, 2, RGBColor(0, 123, 255))], "Python基础"),
    ("⚡ 理论力学", [(1, 2, GREEN)], "专业入门"),
    ("📊 数学建模", [(1, 1, YELLOW)], ""),
    ("🔧 材料力学", [(2, 1, GREEN)], ""),
    ("🔬 进实验室", [(2, 2, PURPLE)], "科研启蒙"),
    ("📐 弹性力学", [(3, 1, GREEN)], ""),
    ("🖥️ 有限元", [(3, 1, RGBColor(0, 123, 255))], "ANSYS"),
    ("🎯 专业深化", [(4, 2, PURPLE)], "CFD/机器人"),
    ("💼 实习/科研", [(4, 2, ORANGE)], ""),
    ("🎓 毕业", [(5, 2, RED)], "保研/秋招")
]

for i, (label, bars, note) in enumerate(entries):
    y = Inches(1.85 + i * 0.42)
    txBox = slide.shapes.add_textbox(Inches(0.2), y, Inches(1.5), Inches(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    for start, length, color in bars:
        x = xs[start] + Inches(0.5)
        w = Inches(length * 1.6 - 0.1)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y + Inches(0.08), w, Inches(0.22))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

# === SLIDE 8: 结语 + 人生建议 ===
slide = add_slide()
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
shape.fill.solid()
shape.fill.fore_color.rgb = BLUE
shape.line.fill.background()

# Quote
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '"力学教会我们的，不是如何计算一根梁的弯矩，'
p.font.size = Pt(22)
p.font.italic = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "而是如何在复杂世界中建立模型、简化问题、逼近真理。"
p2.font.size = Pt(22)
p2.font.italic = True
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER

# Section title
txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12.333), Inches(0.8))
p3 = txBox2.text_frame.paragraphs[0]
p3.text = "人生的四点建议"
p3.font.size = Pt(36)
p3.font.bold = True
p3.font.color.rgb = ORANGE
p3.alignment = PP_ALIGN.CENTER

# Four items with large emoji icons
items = [
    ("📖 英语", ORANGE),
    ("🤖 AI", RGBColor(255, 200, 0)),
    ("💰 定投", RGBColor(100, 200, 100)),
    ("💪 健身", RGBColor(200, 100, 200))
]
for i, (item, color) in enumerate(items):
    x = Inches(1.3 + i * 2.8)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.5), Inches(2.2), Inches(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(x, Inches(4.2), Inches(2.2), Inches(1))
    p = txBox.text_frame.paragraphs[0]
    p.text = item
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# Three action items
add_icon_box(slide, "🚀", "立即行动", Inches(0.5), Inches(6), 0.4, ORANGE)
actions = [
    "找1位学长学姐喝奶茶，聊10分钟科研经验",
    "今晚B站收藏3个Python力学仿真教程",
    '明天高数课坐第一排，课后问老师1个"为什么"'
]
for i, action in enumerate(actions):
    txBox = slide.shapes.add_textbox(Inches(5.5), Inches(5.8 + i * 0.45), Inches(7.5), Inches(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = f"{i+1}. {action}"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE

# Save
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "output.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")